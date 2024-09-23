import glob
import matplotlib.pyplot as plt
import numpy as np
import os

from pptx import Presentation

def main():

    #axial level for output
    jz = 8
    #number of fuel radial nodes
    nf = 22
    #number of clad radial nodes
    nc = 3

    times, timed, fggenv, fgrelv, fgrelp, gpres, z, tfin, tfout, pfc, rfo, rci, bup, qv, gap, hgap, tcin, tcout, esw, et, sigh, sigr, sigz = [], [], [], [], [], [], [], [], [], [], [], [], [], [], [], [], [], [], [], [], [], [], []
    
    #list of outfrd files (could be not consequent)
    flist = sorted(glob.glob('out*'), key=os.path.getmtime)
    for file in flist:
        with open(file, 'r') as f:
            for line in f:
                if line[:8] == 'time (s)' : times.append(float(line[27:]))
                if line[:8] == 'time (d)' : timed.append(float(line[27:]))
                if line[:5] == 'fggen' : fggenv.append(float(line[27:]))
                if line[:11] == 'fgrel (cm3)' : fgrelv.append(float(line[27:]))
                if line[:9] == 'fgrel (%)' : fgrelp.append(float(line[27:]))
                if line[:5] == 'gpres' : gpres.append(float(line[27:]))
                if line[:2] == 'z ' : z = [float(s) for s in line[27:].split()]
                if line[:4] == 'tfin' : tfin.append([float(s) for s in line[27:].split()])
                if line[:5] == 'tfout' : tfout.append([float(s) for s in line[27:].split()])
                if line[:4] == 'tcin' : tcin.append([float(s) for s in line[27:].split()])
                if line[:5] == 'tcout' : tcout.append([float(s) for s in line[27:].split()])
                if line[:3] == 'pfc' : pfc.append([float(s) for s in line[27:].split()])
                if line[:3] == 'rfo' : rfo.append([float(s) for s in line[27:].split()])
                if line[:3] == 'rci' : rci.append([float(s) for s in line[27:].split()])
                if line[:3] == 'bup' : bup.append([float(s) for s in line[27:].split()])
                if line[:2] == 'qv' : qv.append([float(s) for s in line[27:].split()])
                if line[:4] == 'gap ' : gap.append([float(s) for s in line[27:].split()])
                if line[:5] == 'hgap ' : hgap.append([float(s) for s in line[27:].split()])
                if line[:9] == 'eps swell' and int(line[24:27]) == jz : esw.append([float(s) for s in line[27:].split()])
                if line[:6] == 'eps th' and int(line[24:27]) == jz : et.append([float(s) for s in line[27:].split()])
                if line[:5] == 'sig h' and int(line[24:27]) == jz : sigh.append([float(s) for s in line[27:].split()])
                if line[:5] == 'sig r' and int(line[24:27]) == jz : sigr.append([float(s) for s in line[27:].split()])
                if line[:5] == 'sig z' and int(line[24:27]) == jz : sigz.append([float(s) for s in line[27:].split()])
    prs = Presentation()

    #fission gas generation and release
    plt.figure()
    plt.plot(timed, fggenv, label='Generated')
    plt.plot(timed, fgrelv, label='Released')
    plt.xlabel("Time (d)")
    plt.ylabel("Fission gas volume (cm3)")
    plt.grid()
    plt.legend(loc="upper left")
    plt.savefig('plot.png', dpi=300, bbox_inches='tight')
    add_slide(prs)
    plt.close()
    
    #gas pressure
    plt.figure()
    plt.plot(timed, gpres)
    plt.xlabel("Time (d)")
    plt.ylabel("Gas pressure (MPa)")
    plt.grid()
    plt.savefig('plot.png', dpi=300, bbox_inches='tight')
    add_slide(prs)
    plt.close()

    #fuel temperature
    plt.figure()
    tfin = np.transpose(np.array(tfin))
    tfout = np.transpose(np.array(tfout))
    plt.plot(timed, tfin[jz], color='red', label='Fuel center')
    plt.plot(timed, tfout[jz], color='blue', label='Fuel surface')
    plt.xlabel("Time (d)")
    plt.ylabel("Fuel temperature (C) at " + str(z[jz]*1000.0) + " mm")
    plt.grid()
    plt.legend()
    plt.savefig('plot.png', dpi=300, bbox_inches='tight')
    add_slide(prs)
    plt.close()

    #clad temperature
    plt.figure()
    tcin = np.transpose(np.array(tcin))
    tcout = np.transpose(np.array(tcout))
    plt.plot(timed, tcin[jz], color='red', label='Clad center')
    plt.plot(timed, tcout[jz], color='blue', label='Clad surface')
    plt.xlabel("Time (d)")
    plt.ylabel("Clad temperature (C) at " + str(z[jz]*1000.0) + " mm")
    plt.grid()
    plt.legend()
    plt.savefig('plot.png', dpi=300, bbox_inches='tight')
    add_slide(prs)
    plt.close()

    #contact pressure
    plt.figure()
    pfc = np.transpose(np.array(pfc))
    plt.plot(timed, pfc[jz], color='red')
    plt.xlabel("Time (d)")
    plt.ylabel("Contact pressure (MPa) at " + str(z[jz]*1000.0) + " mm")
    plt.grid()
    plt.savefig('plot.png', dpi=300, bbox_inches='tight')
    add_slide(prs)
    plt.close()

    #radii
    plt.figure()
    rfo = np.transpose(np.array(rfo))
    rci = np.transpose(np.array(rci))
    plt.plot(timed, rfo[jz]*1000, color='red', label='Fuel surface')
    plt.plot(timed, rci[jz]*1000, color='blue', label='Clad surface')
    plt.xlabel("Time (d)")
    plt.ylabel("Fuel outer and clad inner radii (mm) at " + str(z[jz]*1000.0) + " mm")
    plt.grid()
    plt.legend()
    plt.savefig('plot.png', dpi=300, bbox_inches='tight')
    add_slide(prs)
    plt.close()

    #power density
    plt.figure()
    qv = np.transpose(np.array(qv))
    plt.plot(timed, qv[jz], color='red')
    plt.xlabel("Time (d)")
    plt.ylabel("Power density (W/m3) at " + str(z[jz]*1000.0) + " mm")
    plt.grid()
    plt.savefig('plot.png', dpi=300, bbox_inches='tight')
    add_slide(prs)
    plt.close()

    #gap
    plt.figure()
    gap = np.transpose(np.array(gap))
    plt.plot(timed, gap[jz], color='red')
    plt.xlabel("Time (d)")
    plt.ylabel("Gap (m) at " + str(z[jz]*1000.0) + " mm")
    plt.grid()
    plt.savefig('plot.png', dpi=300, bbox_inches='tight')
    add_slide(prs)
    plt.close()

    #hgap
    plt.figure()
    hgap = np.transpose(np.array(hgap))
    plt.plot(timed, hgap[jz], color='red')
    plt.xlabel("Time (d)")
    plt.ylabel("Gap conductance (W/m2K) at " + str(z[jz]*1000.0) + " mm")
    plt.grid()
    plt.savefig('plot.png', dpi=300, bbox_inches='tight')
    add_slide(prs)
    plt.close()

    #esw
    plt.figure()
    esw = np.transpose(np.array(esw))
    plt.plot(timed, esw[0], color='red', label='Fuel center')
    plt.plot(timed, esw[nf-1], color='blue', label='Fuel surface')
    plt.xlabel("Time (d)")
    plt.ylabel("Volume swelling strain (%) at " + str(z[jz]*1000.0) + " mm")
    plt.grid()
    plt.legend()
    plt.savefig('plot.png', dpi=300, bbox_inches='tight')
    add_slide(prs)
    plt.close()

    #et
    plt.figure()
    et = np.transpose(np.array(et))
    plt.plot(timed, et[0], color='red', label='Fuel center')
    plt.plot(timed, et[nf-1], color='blue', label='Fuel surface')
    plt.xlabel("Time (d)")
    plt.ylabel("Thermal expansion strain (%) at " + str(z[jz]*1000.0) + " mm")
    plt.grid()
    plt.legend()
    plt.savefig('plot.png', dpi=300, bbox_inches='tight')
    add_slide(prs)
    plt.close()

    #sigh
    plt.figure()
    sigh = np.transpose(np.array(sigh))
    plt.plot(timed, sigh[0], color='red', label='Fuel center')
    plt.plot(timed, sigh[nf-1], color='blue', label='Fuel surface')
    plt.plot(timed, sigh[nf], color='black', label='Clad inner surface')
    plt.plot(timed, sigh[nf+nc-1], color='green', label='Clad outer surface')
    plt.xlabel("Time (d)")
    plt.ylabel("Hoop stress (MPa) at " + str(z[jz]*1000.0) + " mm")
    plt.grid()
    plt.legend()
    plt.savefig('plot.png', dpi=300, bbox_inches='tight')
    add_slide(prs)
    plt.close()

    #sigr
    plt.figure()
    sigr = np.transpose(np.array(sigr))
    plt.plot(timed, sigr[0], color='red', label='Fuel center')
    plt.plot(timed, sigr[nf-1], color='blue', label='Fuel surface')
    plt.plot(timed, sigr[nf], color='black', label='Clad inner surface')
    plt.plot(timed, sigr[nf+nc-1], color='green', label='Clad outer surface')
    plt.xlabel("Time (d)")
    plt.ylabel("Radial stress (MPa) at " + str(z[jz]*1000.0) + " mm")
    plt.grid()
    plt.legend()
    plt.savefig('plot.png', dpi=300, bbox_inches='tight')
    add_slide(prs)
    plt.close()

    #sigz
    plt.figure()
    sigz = np.transpose(np.array(sigz))
    plt.plot(timed, sigz[0], color='red', label='Fuel center')
    plt.plot(timed, sigz[nf-1], color='blue', label='Fuel surface')
    plt.plot(timed, sigz[nf], color='black', label='Clad inner surface')
    plt.plot(timed, sigz[nf+nc-1], color='green', label='Clad outer surface')
    plt.xlabel("Time (d)")
    plt.ylabel("Axial stress (MPa) at " + str(z[jz]*1000.0) + " mm")
    plt.grid()
    plt.legend()
    plt.savefig('plot.png', dpi=300, bbox_inches='tight')
    add_slide(prs)
    plt.close()

    prs.save('plots-time.pptx')
    
    prs = Presentation()

    bup = np.transpose(np.array(bup))

    #fission gas generation and release
    plt.figure()
    plt.plot(bup[jz], fggenv, label='Generated')
    plt.plot(bup[jz], fgrelv, label='Released')
    plt.xlabel("Burnup (MWd/kg)")
    plt.ylabel("Fission gas volume (cm3)")
    plt.grid()
    plt.legend(loc="upper left")
    plt.savefig('plot.png', dpi=300, bbox_inches='tight')
    add_slide(prs)
    plt.close()
    
    #gas pressure
    plt.figure()
    plt.plot(bup[jz], gpres)
    plt.xlabel("Burnup (MWd/kg)")
    plt.ylabel("Gas pressure (MPa)")
    plt.grid()
    plt.savefig('plot.png', dpi=300, bbox_inches='tight')
    add_slide(prs)
    plt.close()

    #fuel temperature
    plt.figure()
    plt.plot(bup[jz], tfin[jz], color='red', label='Fuel center')
    plt.plot(bup[jz], tfout[jz], color='blue', label='Fuel surface')
    plt.xlabel("Burnup (MWd/kg)")
    plt.ylabel("Fuel temperature (C) at " + str(z[jz]*1000.0) + " mm")
    plt.grid()
    plt.legend()
    plt.savefig('plot.png', dpi=300, bbox_inches='tight')
    add_slide(prs)
    plt.close()

    #clad temperature
    plt.figure()
    plt.plot(bup[jz], tcin[jz], color='red', label='Clad center')
    plt.plot(bup[jz], tcout[jz], color='blue', label='Clad surface')
    plt.xlabel("Burnup (MWd/kg)")
    plt.ylabel("Clad temperature (C) at " + str(z[jz]*1000.0) + " mm")
    plt.grid()
    plt.legend()
    plt.savefig('plot.png', dpi=300, bbox_inches='tight')
    add_slide(prs)
    plt.close()

    #contact pressure
    plt.figure()
    plt.plot(bup[jz], pfc[jz], color='red')
    plt.xlabel("Burnup (MWd/kg)")
    plt.ylabel("Contact pressure (MPa) at " + str(z[jz]*1000.0) + " mm")
    plt.grid()
    plt.savefig('plot.png', dpi=300, bbox_inches='tight')
    add_slide(prs)
    plt.close()

    #radii
    plt.figure()
    plt.plot(bup[jz], rfo[jz]*1000, color='red', label='Fuel surface')
    plt.plot(bup[jz], rci[jz]*1000, color='blue', label='Clad surface')
    plt.xlabel("Burnup (MWd/kg)")
    plt.ylabel("Fuel outer and clad inner radii (mm) at " + str(z[jz]*1000.0) + " mm")
    plt.grid()
    plt.legend()
    plt.savefig('plot.png', dpi=300, bbox_inches='tight')
    add_slide(prs)
    plt.close()

    #power density
    plt.figure()
    plt.plot(bup[jz], qv[jz], color='red')
    plt.xlabel("Burnup (MWd)")
    plt.ylabel("Power density (W/m3) at " + str(z[jz]*1000.0) + " mm")
    plt.grid()
    plt.savefig('plot.png', dpi=300, bbox_inches='tight')
    add_slide(prs)
    plt.close()

    #gap
    plt.figure()
    plt.plot(bup[jz], gap[jz], color='red')
    plt.xlabel("Burnup (MWd)")
    plt.ylabel("Gap (m) at " + str(z[jz]*1000.0) + " mm")
    plt.grid()
    plt.savefig('plot.png', dpi=300, bbox_inches='tight')
    add_slide(prs)
    plt.close()

    #hgap
    plt.figure()
    plt.plot(bup[jz], hgap[jz], color='red')
    plt.xlabel("Burnup (MWd)")
    plt.ylabel("Gap conductance (W/m2K) at " + str(z[jz]*1000.0) + " mm")
    plt.grid()
    plt.savefig('plot.png', dpi=300, bbox_inches='tight')
    add_slide(prs)
    plt.close()

    #esw
    plt.figure()
    plt.plot(bup[jz], esw[0], color='red', label='Fuel center')
    plt.plot(bup[jz], esw[nf-1], color='blue', label='Fuel surface')
    plt.xlabel("Burnup (MWd)")
    plt.ylabel("Volume swelling strain (%) at " + str(z[jz]*1000.0) + " mm")
    plt.grid()
    plt.legend()
    plt.savefig('plot.png', dpi=300, bbox_inches='tight')
    add_slide(prs)
    plt.close()

    #et
    plt.figure()
    plt.plot(bup[jz], et[0], color='red', label='Fuel center')
    plt.plot(bup[jz], et[nf-1], color='blue', label='Fuel surface')
    plt.xlabel("Burnup (MWd)")
    plt.ylabel("Thermal expansion strain (%) at " + str(z[jz]*1000.0) + " mm")
    plt.grid()
    plt.legend()
    plt.savefig('plot.png', dpi=300, bbox_inches='tight')
    add_slide(prs)
    plt.close()

    #sigh
    plt.figure()
    plt.plot(bup[jz], sigh[0], color='red', label='Fuel center')
    plt.plot(bup[jz], sigh[nf-1], color='blue', label='Fuel surface')
    plt.plot(bup[jz], sigh[nf], color='black', label='Clad inner surface')
    plt.plot(bup[jz], sigh[nf+nc-1], color='green', label='Clad outer surface')
    plt.xlabel("Burnup (MWd)")
    plt.ylabel("Hoop stress (MPa) at " + str(z[jz]*1000.0) + " mm")
    plt.grid()
    plt.legend()
    plt.savefig('plot.png', dpi=300, bbox_inches='tight')
    add_slide(prs)
    plt.close()

    #sigr
    plt.figure()
    plt.plot(bup[jz], sigr[0], color='red', label='Fuel center')
    plt.plot(bup[jz], sigr[nf-1], color='blue', label='Fuel surface')
    plt.plot(bup[jz], sigr[nf], color='black', label='Clad inner surface')
    plt.plot(bup[jz], sigr[nf+nc-1], color='green', label='Clad outer surface')
    plt.xlabel("Burnup (MWd)")
    plt.ylabel("Radial stress (MPa) at " + str(z[jz]*1000.0) + " mm")
    plt.grid()
    plt.legend()
    plt.savefig('plot.png', dpi=300, bbox_inches='tight')
    add_slide(prs)
    plt.close()

    #sigz
    plt.figure()
    plt.plot(bup[jz], sigz[0], color='red', label='Fuel center')
    plt.plot(bup[jz], sigz[nf-1], color='blue', label='Fuel surface')
    plt.plot(bup[jz], sigz[nf], color='black', label='Clad inner surface')
    plt.plot(bup[jz], sigz[nf+nc-1], color='green', label='Clad outer surface')
    plt.xlabel("Burnup (MWd)")
    plt.ylabel("Axial stress (MPa) at " + str(z[jz]*1000.0) + " mm")
    plt.grid()
    plt.legend()
    plt.savefig('plot.png', dpi=300, bbox_inches='tight')
    add_slide(prs)
    plt.close()
    
    prs.save('plots-bup.pptx')

def add_slide(prs):

    img_width = 5047488
    img_height = 3776472

    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    for placeholder in slide.shapes.placeholders:
        sp = placeholder._sp
        sp.getparent().remove(sp)
    left = (prs.slide_width - img_width) / 2
    top = (prs.slide_height - img_height) / 2
    img = slide.shapes.add_picture('./plot.png',left,top)
    
if __name__=="__main__":
    main()